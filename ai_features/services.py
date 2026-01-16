"""
خدمات الذكاء الاصطناعي
S-ACM - Smart Academic Content Management System
"""

import os
import json
from pathlib import Path
from django.conf import settings

# OpenAI-compatible API (يدعم Gemini)
try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

# PDF Processing
try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Word Processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint Processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class GeminiService:
    """خدمة Google Gemini للذكاء الاصطناعي عبر OpenAI-compatible API"""
    
    def __init__(self):
        if OPENAI_AVAILABLE:
            self.client = OpenAI()  # يستخدم المتغيرات البيئية تلقائياً
            self.model = "gemini-2.5-flash"
        else:
            self.client = None
            self.model = None
    
    def is_available(self):
        """التحقق من توفر الخدمة"""
        return self.client is not None
    
    def extract_text_from_file(self, file_obj):
        """استخراج النص من الملف"""
        if file_obj.content_type == 'external_link':
            return None  # لا يمكن استخراج النص من الروابط الخارجية
        
        if not file_obj.local_file:
            return None
        
        file_path = Path(file_obj.local_file.path)
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pdf':
                return self._extract_from_pdf(file_path)
            elif extension == '.docx':
                return self._extract_from_docx(file_path)
            elif extension == '.pptx':
                return self._extract_from_pptx(file_path)
            elif extension in ['.txt', '.md']:
                return self._extract_from_text(file_path)
            else:
                return None
        except Exception as e:
            print(f"Error extracting text: {e}")
            return None
    
    def _extract_from_pdf(self, file_path):
        """استخراج النص من PDF"""
        if not PDF_AVAILABLE:
            return None
        
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    
    def _extract_from_docx(self, file_path):
        """استخراج النص من Word"""
        if not DOCX_AVAILABLE:
            return None
        
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    
    def _extract_from_pptx(self, file_path):
        """استخراج النص من PowerPoint"""
        if not PPTX_AVAILABLE:
            return None
        
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    
    def _extract_from_text(self, file_path):
        """استخراج النص من ملف نصي"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read().strip()
    
    def generate_summary(self, text, max_length=500):
        """توليد تلخيص للنص"""
        if not self.is_available():
            return self._fallback_summary(text, max_length)
        
        prompt = f"""
        أنت مساعد أكاديمي متخصص في تلخيص المحتوى التعليمي.
        قم بتلخيص النص التالي بشكل مختصر ومفيد باللغة العربية.
        ركز على النقاط الرئيسية والمفاهيم الأساسية.
        
        النص:
        {text[:10000]}
        
        التلخيص:
        """
        
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "أنت مساعد أكاديمي متخصص في تلخيص المحتوى التعليمي باللغة العربية."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=max_length,
                temperature=0.3
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"AI error: {e}")
            return self._fallback_summary(text, max_length)
    
    def _fallback_summary(self, text, max_length):
        """تلخيص بسيط في حالة عدم توفر الخدمة"""
        sentences = text.split('.')
        summary = ""
        for sentence in sentences:
            if len(summary) + len(sentence) < max_length:
                summary += sentence.strip() + ". "
            else:
                break
        return summary.strip() or text[:max_length] + "..."
    
    def generate_questions(self, text, question_type='mixed', num_questions=5):
        """توليد أسئلة من النص"""
        if not self.is_available():
            return self._fallback_questions(text, num_questions)
        
        type_instruction = {
            'mcq': 'أسئلة اختيار من متعدد فقط',
            'true_false': 'أسئلة صح أو خطأ فقط',
            'short_answer': 'أسئلة إجابة قصيرة فقط',
            'mixed': 'مزيج من أنواع الأسئلة المختلفة'
        }.get(question_type, 'مزيج من أنواع الأسئلة')
        
        prompt = f"""
        أنت مدرس متخصص في إنشاء أسئلة اختبارية.
        قم بإنشاء {num_questions} سؤال من النص التالي.
        نوع الأسئلة المطلوب: {type_instruction}
        
        أرجع الإجابة بصيغة JSON كالتالي:
        [
            {{
                "type": "mcq" أو "true_false" أو "short_answer",
                "question": "نص السؤال",
                "options": ["خيار1", "خيار2", "خيار3", "خيار4"] (للاختيار من متعدد فقط),
                "answer": "الإجابة الصحيحة",
                "explanation": "شرح مختصر للإجابة"
            }}
        ]
        
        النص:
        {text[:8000]}
        
        الأسئلة (JSON فقط):
        """
        
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "أنت مدرس متخصص في إنشاء أسئلة اختبارية تعليمية باللغة العربية. قدم الإجابة بصيغة JSON فقط."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000,
                temperature=0.5
            )
            
            response_text = response.choices[0].message.content.strip()
            
            # محاولة استخراج JSON
            if '```json' in response_text:
                response_text = response_text.split('```json')[1].split('```')[0]
            elif '```' in response_text:
                response_text = response_text.split('```')[1].split('```')[0]
            
            questions = json.loads(response_text)
            return questions
        except Exception as e:
            print(f"AI error: {e}")
            return self._fallback_questions(text, num_questions)
    
    def _fallback_questions(self, text, num_questions):
        """أسئلة بسيطة في حالة عدم توفر الخدمة"""
        return [
            {
                'type': 'short_answer',
                'question': 'ما هي الفكرة الرئيسية في هذا النص؟',
                'answer': 'راجع النص للإجابة',
                'explanation': 'هذا سؤال تلقائي'
            }
        ]
    
    def ask_document(self, text, question):
        """الإجابة على سؤال من سياق المستند"""
        if not self.is_available():
            return "عذراً، خدمة الذكاء الاصطناعي غير متاحة حالياً."
        
        prompt = f"""
        أنت مساعد أكاديمي ذكي. أجب على السؤال التالي بناءً على المحتوى المقدم فقط.
        إذا لم تجد الإجابة في المحتوى، قل ذلك بوضوح.
        
        المحتوى:
        {text[:10000]}
        
        السؤال: {question}
        
        الإجابة:
        """
        
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": "أنت مساعد أكاديمي يجيب على الأسئلة بناءً على محتوى المستندات المقدمة. أجب باللغة العربية بشكل واضح ومفيد."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=500,
                temperature=0.3
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"AI error: {e}")
            return "عذراً، حدث خطأ أثناء معالجة سؤالك. يرجى المحاولة مرة أخرى."
